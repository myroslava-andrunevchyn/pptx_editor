"""
Fills in necessary fields in pptx-template with data from mysql database
"""

import mysql.connector
from pptx import Presentation

'''
Connection to Database Server
'''
db = mysql.connector.connect(
    host="127.0.0.1",
    user="root",
    password="135790Server!",
    database="clients"
)
if db.is_connected():
    print("Database Connected")


def get_clients(database):
    """Selects all clients ids from database and returns the list of IDs for detailed search."""
    cursor = database.cursor()
    sql = 'SELECT accounts.id FROM accounts'
    cursor.execute(sql)
    clients_list = cursor.fetchall()
    return clients_list


def get_client_data(database):
    """Iterating trough clients_list retrieves data and returns tuple of 2 dictionaries: 1 - string data,
    2 - graphic data."""
    cursor = database.cursor()
    clients_ids = get_clients(database)

    tags_string = ('{{client_name}}', '{{budget}}', '{{contact}}', '{{date_est}}', '{{mail}}')   # keys for dictionary with clients string data
    tags_pictures = ('{{logo}}', '{{city_picture}}', '{{diagram}}')   # keys for dictionary with clients graphic data

    sql = '''SELECT accounts.client_name, reports.budget, reports.contact, reports.date_est, reports.mail 
                FROM accounts inner join reports on accounts.id = reports.account_id where account_id=%s'''   # select for string data per client
    sql2 = '''SELECT reports.logo, reports.city_picture, reports.diagram 
                FROM accounts inner join reports on accounts.id = reports.account_id where account_id=%s'''   # select for graphic data per client

    for client_id in clients_ids:
        cursor.execute(sql, client_id)
        client_data_for_strings = cursor.fetchall()
        client_data_list_for_strings = dict(zip(tags_string, client_data_for_strings[0]))
        cursor.execute(sql2, client_id)
        client_data_for_pictures = cursor.fetchall()
        client_data_list_for_pictures = dict(zip(tags_pictures, client_data_for_pictures[0]))
        yield client_data_list_for_strings, client_data_list_for_pictures   # generates dictionary with string data and dictionary with graphic data per client


def create_pptx_report(database):
    """Inserts data from database and saves separate pptx file with company name."""
    client_data = get_client_data(database)   # tuple of 2 dictionaries with data per company

    for company_data in client_data:
        pptx_template = Presentation(r"Template for Project.pptx")   # pptx template file
        for slide in pptx_template.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:  # search for tags in text to replace with client data
                    for paragraph in shape.text_frame.paragraphs:
                        replace_text(company_data, paragraph)   # replaces tag in template file with corresponding client data
                for key in company_data[0]:   # replaces data in placeholders with corresponding name
                    if shape.name == key:
                        shape.text = company_data[0][key]
            for placeholder in slide.placeholders:   # inserts picture into placeholders with corresponding name
                insert_picture(company_data, placeholder)

        pptx_template.save('%s edited.pptx' % company_data[0]['{{client_name}}'])   # saves edited template as ptx file with a company name


def replace_text(company_data, paragraph):
    """Replaces tags with corresponding string data from database."""
    for key in company_data[0]:
        if paragraph.text.find(key) != -1:
            paragraph.text = paragraph.text.replace(key, str(company_data[0][key]))


def insert_picture(company_data, placeholder):
    """Inserts graphic object into placeholder with corresponding tag-name."""
    for key in company_data[1]:
        if placeholder.name == key:
            placeholder.insert_picture(company_data[1][key])
            break


if __name__ == "__main__":
    create_pptx_report(db)